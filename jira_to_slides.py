# Author: Luke Connolly
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
import configparser
import requests as r
 
config = configparser.ConfigParser()
config.read('config.ini')
 
try:
    RAPID_VIEW_ID = str(config['CONFIG']['RAPID_VIEW_ID'])
    JSESSIONID = str(config['CONFIG']['JSESSION_ID'])
    TEAM_NAME = str(config['CONFIG']['TEAM_NAME'])
    BACKLOG = config['CONFIG'].getboolean('BACKLOG')
    if(RAPID_VIEW_ID == '' or JSESSIONID == '' or TEAM_NAME == ''):
        raise ValueError()
 
except ValueError:
    print('Invalid config.ini file. Please check your config.ini file and try again.')
    exit(1)
 
JIRA_URL = f'https://issues.corp.rapid7.com/rest/greenhopper/1.0/xboard/plan/backlog/data.json?rapidViewId={RAPID_VIEW_ID}&selectedProjectKey=LOG'
COOKIES = {'JSESSIONID': JSESSIONID}
 
print(f'Generating presentation for {TEAM_NAME}')
print('Getting JIRA Data...')
def getSprintData() -> dict:
    res = r.get(JIRA_URL, cookies=COOKIES)
    if res.status_code == 401:
        raise ConnectionError('Could not get JIRA Data. Check your JSESSIONID, it may have expired.')
    return res.json()
 
SPRINT = getSprintData()
EPIC_DATA = SPRINT['entityData']['epics']
ALL_ISSUES = SPRINT['issues']
THIS_SPRINT_ISSUE_IDS = SPRINT['sprints'][0]['issuesIds']
TITLE = f'{TEAM_NAME} Sprint Work'
BUGS = 'Bugs'
PRS = Presentation('template.pptx')
LAYOUT = PRS.slide_layouts[11]
 
class Issue:
    def __init__(self, summary, status, epic=None):
        self.summary = summary
        self.status = status
        self.epic = epic
 
    def __str__(self):
        return f"{self.summary} - {self.status}"
 
class Epic:
    def addIssue(self, issue):
        self.issues.append(issue)
 
    def __init__(self, name):
        self.name = name
        self.issues = []
 
    def __str__(self):
        return self.name
 
    def __len__(self):
        return len(self.issues)
 
    def __lt__(self, other):
        return len(self.issues) < len(other.issues)
 
def getEpicName(epicId:str) -> str:
    return EPIC_DATA[epicId]['epicField']['text']
 
def getIssueStatus(statusId:str) -> str:
    return SPRINT['entityData']['statuses'][statusId]['statusName']
 
def isInThisSprint(issue:object) -> bool:
    return issue['id'] in THIS_SPRINT_ISSUE_IDS
 
def isBug(issue:object) -> bool:
    return 'epicId' not in issue
 
def getTotalFinishedBugs(bugs: list) -> str:
    return len(list(filter(lambda bug : bug.status == 'Done', bugs)))
 
def getBugTitle(bugs: list) -> str:
    closedBugs = getTotalFinishedBugs(bugs)
    return f"Bugs ({closedBugs} closed, {len(bugs) - closedBugs} open ) "
 
def createSlide(title:str):
    slide = PRS.slides.add_slide(LAYOUT)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    return body.text_frame
 
def addEpicToSlide(slide, epic:Epic):
    paragraph = slide.paragraphs[0] if slide.paragraphs[0].text == '' else slide.add_paragraph()
    paragraph.text = epic.name
    paragraph.font.size = Pt(17.5)
    paragraph.font.bold = True
 
    for issue in epic.issues:
        if(BACKLOG and issue.status == 'Backlog'):
            continue
        paragraph = slide.add_paragraph()
        paragraph.text = str(issue)
        paragraph.font.size = Pt(13.5)
        paragraph.level = 1
 
    #add a blank line
    paragraph = slide.add_paragraph()
    paragraph.text = ''
 
epics = {}
bugs = []
 
print('Parsing JIRA Data...')
for issue in filter(isInThisSprint, ALL_ISSUES):
    if isBug(issue):
        bugs.append(Issue(issue['summary'], getIssueStatus(issue['statusId'])))
        continue
 
    issue = Issue(issue['summary'], getIssueStatus(issue['statusId']), getEpicName(issue['epicId']))
    if issue.epic not in epics:
        epics[issue.epic] = Epic(issue.epic)
    epics[issue.epic].addIssue(issue)
 
epicsSortedBySize = sorted(epics.values(), reverse=True)
 
print(f'Generating Presentation for {len(epicsSortedBySize)} Epics, {len(bugs)} Bugs.')
totalIssues = 0
totalEpics = 0
slide = createSlide(TITLE)
for epic in epicsSortedBySize:
    if totalIssues + len(epic) > 10 or (totalEpics > 2 and totalIssues + len(epic) > 5):
        slide = createSlide(TITLE)
        totalIssues = 0
        totalEpics = 0
 
    addEpicToSlide(slide, epic)
    totalIssues += len(epic)
    totalEpics += 1
 
slide = createSlide(BUGS)
slide.text = getBugTitle(bugs)
slide.paragraphs[0].font.size = Pt(17.5)
slide.paragraphs[0].font.bold = True
for bug in bugs:
    paragraph = slide.add_paragraph()
    paragraph.text = str(bug)
    paragraph.font.size = Pt(13.5)
    paragraph.level = 1
 
#Remove Template Slide
xml_slides = PRS.slides._sldIdLst
slides = list(xml_slides)
xml_slides.remove(slides[0])
 
PRS.save('Sprint.pptx')
print('Success! \nSaved to Sprint.pptx')